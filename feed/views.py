"""
Views (представлення) для SchoolNet.

Публічна частина:
    - index          — стрічка завдань з фільтрацією по класу
    - assignment_detail — детальна сторінка/AJAX для модального вікна

Панель вчителя (захищена сесією):
    - teacher_login      — форма входу
    - teacher_logout     — вихід
    - teacher_dashboard  — головна панель (список завдань)
    - assignment_create  — створення нового завдання
    - assignment_edit    — редагування завдання
    - assignment_delete  — видалення завдання
    - assignment_duplicate — дублювання завдання
    - assignment_archive — архівування завдання
    - teacher_profile    — редагування профілю
    - publish_scheduled  — AJAX: перевірка та публікація відкладених
"""

from django.shortcuts import render, redirect, get_object_or_404
from django.urls import reverse
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.models import User
from django.contrib import messages
from django.http import JsonResponse, HttpResponse
from django.utils import timezone
from django.db.models import Q
from django.core.paginator import Paginator, EmptyPage, PageNotAnInteger

from .models import Assignment, AssignmentFile, AssignmentLink, AssignmentYouTubeLink, Teacher, ClassGroup, Subject
from .forms import (
    AssignmentForm, TeacherLoginForm, ClassSelectForm,
    TeacherProfileForm, SubjectForm, ClassGroupForm,
    TeacherCreateForm, PasswordResetForm
)
from .search import search_assignments


# ═══════════════════════════════════════════════════════════════════════════════
# ДОПОМІЖНІ ФУНКЦІЇ
# ═══════════════════════════════════════════════════════════════════════════════

def get_teacher_or_none(request):
    """Повертає профіль вчителя або None якщо не авторизовано."""
    if request.user.is_authenticated:
        try:
            return request.user.teacher_profile
        except Teacher.DoesNotExist:
            pass
    return None


def teacher_required(view_func):
    """Декоратор: перевіряє чи є активна сесія вчителя."""
    def wrapper(request, *args, **kwargs):
        teacher = get_teacher_or_none(request)
        if not teacher:
            messages.warning(request, 'Необхідно увійти як вчитель.')
            return redirect('teacher_login')
        return view_func(request, *args, **kwargs)
    wrapper.__name__ = view_func.__name__
    return wrapper


def auto_archive_expired_assignments():
    """
    Автоматично архівує всі опубліковані завдання, яким понад 14 днів з моменту публікації або розархівації.
    """
    from datetime import timedelta
    cutoff = timezone.now() - timedelta(days=14)
    # Завдання без розархівації — рахуємо від первинної дати публікації / створення
    q_normal = Q(unarchived_at__isnull=True) & (
        Q(published_at__lte=cutoff) | (Q(published_at__isnull=True) & Q(created_at__lte=cutoff))
    )
    # Завдання з розархівацією — рахуємо 14 днів від дати розархівації
    q_unarchived = Q(unarchived_at__isnull=False) & Q(unarchived_at__lte=cutoff)

    expired_count = Assignment.objects.filter(
        status=Assignment.STATUS_PUBLISHED
    ).filter(
        q_normal | q_unarchived
    ).update(
        status=Assignment.STATUS_ARCHIVED
    )
    return expired_count


def get_visible_assignments(class_group_id=None):
    """
    Повертає QuerySet опублікованих завдань.
    - Автоматично архівує завдання старше 14 днів.
    - Автоматично публікує відкладені, якщо час настав.
    """
    # Автоматична архівація завдань, старших за 14 днів
    auto_archive_expired_assignments()

    # Оновлюємо відкладені публікації
    Assignment.objects.filter(
        status=Assignment.STATUS_SCHEDULED,
        scheduled_at__lte=timezone.now()
    ).update(
        status=Assignment.STATUS_PUBLISHED,
        published_at=timezone.now()
    )

    qs = Assignment.objects.filter(
        status=Assignment.STATUS_PUBLISHED
    ).select_related(
        'teacher', 'subject'
    ).prefetch_related(
        'classes', 'files'
    )

    if class_group_id:
        qs = qs.filter(classes__id=class_group_id)

    return qs.order_by('-published_at')


# ═══════════════════════════════════════════════════════════════════════════════
# ПУБЛІЧНА ЧАСТИНА — СТРІЧКА НОВИН
# ═══════════════════════════════════════════════════════════════════════════════

def get_calendar_context(year=None, month=None, selected_date_str=None, class_group_id=None, subject_id=None):
    """
    Формує матрицю днів та кількість завдань для віджета Google Calendar.
    """
    import calendar
    from datetime import timedelta, datetime
    from django.utils import timezone
    from collections import Counter

    today = timezone.localtime(timezone.now()).date()

    if not year or not month:
        if selected_date_str:
            try:
                d = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
                year, month = d.year, d.month
            except Exception:
                year, month = today.year, today.month
        else:
            year, month = today.year, today.month
    else:
        try:
            year = int(year)
            month = int(month)
        except Exception:
            year, month = today.year, today.month

    # Назви місяців українською
    uk_months = [
        '', 'Січень', 'Лютий', 'Березень', 'Квітень', 'Травень', 'Червень',
        'Липень', 'Серпень', 'Вересень', 'Жовтень', 'Листопад', 'Грудень'
    ]

    # Навігація по місяцях
    if month <= 1:
        prev_month, prev_year = 12, year - 1
    else:
        prev_month, prev_year = month - 1, year

    if month >= 12:
        next_month, next_year = 1, year + 1
    else:
        next_month, next_year = month + 1, year

    # Отримуємо всі опубліковані завдання для підрахунку на календарі
    visible_qs = get_visible_assignments(class_group_id)
    if subject_id:
        visible_qs = visible_qs.filter(subject__id=subject_id)

    task_counts = Counter()
    for a in visible_qs:
        if a.published_at:
            p_date = timezone.localtime(a.published_at).date()
            task_counts[p_date.isoformat()] += 1
        if a.unarchived_at:
            u_date = timezone.localtime(a.unarchived_at).date()
            p_date = timezone.localtime(a.published_at).date() if a.published_at else None
            if u_date != p_date:
                task_counts[u_date.isoformat()] += 1

    cal = calendar.Calendar(firstweekday=0)  # Понеділок (0)
    month_weeks = cal.monthdatescalendar(year, month)

    fourteen_days_ago = today - timedelta(days=14)

    weeks_data = []
    for week in month_weeks:
        week_days = []
        for d in week:
            d_str = d.isoformat()
            c = task_counts.get(d_str, 0)
            week_days.append({
                'date': d,
                'day_num': d.day,
                'date_str': d_str,
                'is_current_month': (d.month == month),
                'is_today': (d == today),
                'is_selected': (d_str == selected_date_str),
                'is_in_14_days': (fourteen_days_ago <= d <= today),
                'tasks_count': c,
                'has_tasks': c > 0,
            })
        weeks_data.append(week_days)

    selected_date_display = ""
    if selected_date_str:
        try:
            dt = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            uk_weekdays = ['Понеділок', 'Вівторок', 'Середа', 'Четвер', 'Пʼятниця', 'Субота', 'Неділя']
            w_name = uk_weekdays[dt.weekday()]
            selected_date_display = f"{w_name}, {dt.strftime('%d.%m.%Y')}"
        except Exception:
            pass

    return {
        'year': year,
        'month': month,
        'month_name': uk_months[month],
        'prev_year': prev_year,
        'prev_month': prev_month,
        'next_year': next_year,
        'next_month': next_month,
        'weeks': weeks_data,
        'selected_date': selected_date_str,
        'selected_date_display': selected_date_display,
        'today_str': today.isoformat(),
    }


def _parse_filter_params(request):
    """Допоміжна функція: парсить та очищує параметри фільтрації класу, предмета, дати та пошуку."""
    class_group_id = None
    if 'class' in request.GET:
        c_val = str(request.GET.get('class', '')).strip()
        if c_val in ('', 'all', '0', 'none', 'None'):
            class_group_id = None
            if 'selected_class' in request.session:
                del request.session['selected_class']
        else:
            try:
                class_group_id = int(c_val)
                request.session['selected_class'] = class_group_id
            except (ValueError, TypeError):
                class_group_id = None
                if 'selected_class' in request.session:
                    del request.session['selected_class']
    elif 'clear' in request.GET:
        class_group_id = None
        if 'selected_class' in request.session:
            del request.session['selected_class']
    else:
        class_group_id = request.session.get('selected_class')

    subject_id = None
    if 'subject' in request.GET:
        s_val = str(request.GET.get('subject', '')).strip()
        if s_val not in ('', 'all', '0', 'none', 'None'):
            try:
                subject_id = int(s_val)
            except (ValueError, TypeError):
                subject_id = None

    date_str = None
    if 'date' in request.GET:
        d_val = str(request.GET.get('date', '')).strip()
        if d_val not in ('', 'all', '0', 'none', 'None'):
            try:
                from datetime import datetime
                datetime.strptime(d_val, '%Y-%m-%d')
                date_str = d_val
            except (ValueError, TypeError):
                date_str = None

    query = request.GET.get('q', '').strip()
    return class_group_id, subject_id, date_str, query


def index(request):
    """
    Головна сторінка — стрічка завдань.
    Підтримує фільтрацію по класу, предмету, даті (календар) та нечіткий пошук.
    """
    class_group_id, subject_id, date_str, query = _parse_filter_params(request)

    assignments = get_visible_assignments(class_group_id)

    if subject_id:
        assignments = assignments.filter(subject__id=subject_id)

    if date_str:
        from datetime import datetime
        dt_val = datetime.strptime(date_str, '%Y-%m-%d').date()
        assignments = assignments.filter(
            Q(published_at__date=dt_val) | Q(unarchived_at__date=dt_val)
        )

    if query:
        assignments = search_assignments(assignments, query)

    # Пагінація (10 завдань на сторінку)
    paginator = Paginator(assignments, 10)
    page_number = request.GET.get('page', 1)
    try:
        page_obj = paginator.page(page_number)
    except PageNotAnInteger:
        page_obj = paginator.page(1)
    except EmptyPage:
        page_obj = paginator.page(paginator.num_pages)

    # Дані для UI
    all_classes = ClassGroup.objects.filter(teacher__isnull=False).distinct().order_by('grade', 'letter')
    all_subjects = Subject.objects.filter(teacher__isnull=False).distinct()
    has_multiple_subjects = all_subjects.count() > 1
    selected_class = None
    if class_group_id:
        selected_class = ClassGroup.objects.filter(id=class_group_id).first()

    cal_year = request.GET.get('cal_year')
    cal_month = request.GET.get('cal_month')
    calendar_ctx = get_calendar_context(
        year=cal_year,
        month=cal_month,
        selected_date_str=date_str,
        class_group_id=class_group_id,
        subject_id=subject_id
    )

    context = {
        'assignments': page_obj.object_list,
        'page_obj': page_obj,
        'all_classes': all_classes,
        'all_subjects': all_subjects,
        'has_multiple_subjects': has_multiple_subjects,
        'selected_class': selected_class,
        'selected_class_id': class_group_id,
        'selected_subject_id': subject_id,
        'selected_date': date_str,
        'selected_date_display': calendar_ctx.get('selected_date_display', ''),
        'calendar': calendar_ctx,
        'query': query,
        'form': ClassSelectForm(initial={'class_group': class_group_id}),
    }
    return render(request, 'feed/index.html', context)



def assignment_detail(request, pk):
    """
    Детальний перегляд завдання.
    Якщо AJAX — повертає JSON для модального вікна.
    Інакше — повна сторінка.
    """
    assignment = get_object_or_404(
        Assignment.objects.select_related('teacher', 'subject').prefetch_related('classes', 'files', 'additional_links', 'youtube_links'),
        pk=pk,
        status=Assignment.STATUS_PUBLISHED
    )

    # Фіксація перегляду завдання (1 раз на годину з 1 комп'ютера, перегляди вчителя не рахуються)
    assignment.record_view(request)

    is_teacher = bool(
        request.user.is_authenticated and
        (hasattr(request.user, 'teacher_profile') or request.user.is_superuser)
    )

    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        # AJAX-відповідь для модального вікна
        files_data = []
        for f in assignment.files.all():
            files_data.append({
                'id': f.id,
                'name': f.original_name,
                'url': f.file.url,
                'type': f.get_file_type(),
                'icon': f.get_file_icon(),
                'size': f.get_size_display(),
                'viewable': f.is_browser_viewable(),
                'extension': f.get_extension(),
            })

        # Формуємо рядок класів
        classes_str = ', '.join(
            str(c) for c in assignment.classes.all()
        )

        # Додаткові посилання
        extra_links = [
            {'url': lnk.url, 'label': lnk.label or lnk.url}
            for lnk in assignment.additional_links.all()
        ]

        data = {
            'id': assignment.id,
            'title': assignment.title,
            'description': assignment.description,
            'teacher': assignment.teacher.full_name,
            'teacher_avatar_url': assignment.teacher.avatar_image.url if assignment.teacher.avatar_image else '',
            'teacher_avatar_color': assignment.teacher.avatar_color,
            'subject': str(assignment.subject) if assignment.subject else '',
            'subject_color': assignment.subject.color if assignment.subject else '#6366f1',
            'classes': classes_str,
            'is_individual': assignment.is_individual,
            'student_name': assignment.student_name,
            'student_display': assignment.student_display,
            'views_count': assignment.views_count,
            'is_teacher': is_teacher,
            'published_at': assignment.published_at.strftime('%d.%m.%Y %H:%M') if assignment.published_at else '',
            'relative_published_at': assignment.relative_published_display,
            'is_today': assignment.is_published_today,
            'is_yesterday': assignment.is_published_yesterday,
            'days_ago': assignment.days_ago,
            'unarchived_at': assignment.unarchived_at.strftime('%d.%m.%Y %H:%M') if assignment.unarchived_at else '',
            'unarchived_display': assignment.unarchived_display,
            'due_date': assignment.due_date.strftime('%d.%m.%Y') if assignment.due_date else '',
            'link_url': assignment.link_url,
            'link_label': assignment.link_label or assignment.link_url,
            'youtube_url': assignment.youtube_embed_url,
            'youtube_watch_url': assignment.youtube_watch_url,
            'youtube_videos': assignment.all_youtube_videos,
            'extra_links': extra_links,
            'files': files_data,
            'can_edit': bool(
                is_teacher and
                (request.user.is_superuser or (hasattr(request.user, 'teacher_profile') and request.user.teacher_profile == assignment.teacher))
            ),
        }
        return JsonResponse(data)

    context = {
        'assignment': assignment,
        'is_teacher': is_teacher,
    }
    return render(request, 'feed/assignment_detail.html', context)


def clear_class_filter(request):
    """Скидає фільтр класу з сесії."""
    if 'selected_class' in request.session:
        del request.session['selected_class']
    return redirect('index')


def feed_fragment(request):
    """
    AJAX: повертає лише HTML фрагмент стрічки завдань (картки + пагінація).
    Використовується при клінці на фільтр класу/предмету/дати без перезавантаження сторінки.
    """
    class_group_id, subject_id, date_str, query = _parse_filter_params(request)
    page_number = request.GET.get('page', 1)

    assignments = get_visible_assignments(class_group_id)

    if subject_id:
        assignments = assignments.filter(subject__id=subject_id)

    if date_str:
        from datetime import datetime
        dt_val = datetime.strptime(date_str, '%Y-%m-%d').date()
        assignments = assignments.filter(
            Q(published_at__date=dt_val) | Q(unarchived_at__date=dt_val)
        )

    if query:
        assignments = search_assignments(assignments, query)

    paginator = Paginator(assignments, 10)
    try:
        page_obj = paginator.page(page_number)
    except PageNotAnInteger:
        page_obj = paginator.page(1)
    except EmptyPage:
        page_obj = paginator.page(paginator.num_pages)

    all_subjects = Subject.objects.filter(teacher__isnull=False).distinct()
    has_multiple_subjects = all_subjects.count() > 1
    selected_class = ClassGroup.objects.filter(id=class_group_id).first() if class_group_id else None

    selected_date_display = ""
    if date_str:
        try:
            from datetime import datetime
            dt = datetime.strptime(date_str, '%Y-%m-%d').date()
            uk_weekdays = ['Понеділок', 'Вівторок', 'Середа', 'Четвер', 'Пʼятниця', 'Субота', 'Неділя']
            w_name = uk_weekdays[dt.weekday()]
            selected_date_display = f"{w_name}, {dt.strftime('%d.%m.%Y')}"
        except Exception:
            pass

    context = {
        'assignments': page_obj.object_list,
        'page_obj': page_obj,
        'selected_class': selected_class,
        'selected_class_id': class_group_id,
        'selected_subject_id': subject_id,
        'selected_date': date_str,
        'selected_date_display': selected_date_display,
        'has_multiple_subjects': has_multiple_subjects,
        'query': query,
    }
    return render(request, 'feed/feed_fragment.html', context)


def calendar_fragment(request):
    """
    AJAX: повертає HTML-фрагмент календаря при зміні місяця або фільтрів.
    """
    class_group_id, subject_id, date_str, _ = _parse_filter_params(request)
    cal_year = request.GET.get('cal_year')
    cal_month = request.GET.get('cal_month')

    calendar_ctx = get_calendar_context(
        year=cal_year,
        month=cal_month,
        selected_date_str=date_str,
        class_group_id=class_group_id,
        subject_id=subject_id
    )
    return render(request, 'feed/calendar_widget.html', {'calendar': calendar_ctx})


def feed_check_updates(request):
    """
    Надлегка перевірка наявності нових або змінених завдань для стрічки.
    Виконує 1 швидкий запит агрегації в БД (<1 мс), не навантажуючи сервер.
    """
    class_group_id, subject_id, date_str, query = _parse_filter_params(request)

    assignments = get_visible_assignments(class_group_id)
    if subject_id:
        assignments = assignments.filter(subject__id=subject_id)
    if date_str:
        from datetime import datetime
        dt_val = datetime.strptime(date_str, '%Y-%m-%d').date()
        assignments = assignments.filter(
            Q(published_at__date=dt_val) | Q(unarchived_at__date=dt_val)
        )
    if query:
        assignments = assignments.filter(
            Q(title__icontains=query) |
            Q(description__icontains=query) |
            Q(teacher__full_name__icontains=query)
        )

    from django.db.models import Max, Count
    agg = assignments.aggregate(
        total_count=Count('id'),
        max_id=Max('id'),
        max_pub=Max('published_at')
    )

    pub_ts = int(agg['max_pub'].timestamp()) if agg['max_pub'] else 0
    fingerprint = f"{agg['total_count']}-{agg['max_id'] or 0}-{pub_ts}"

    return JsonResponse({
        'fingerprint': fingerprint,
        'count': agg['total_count'] or 0,
    })


# ═══════════════════════════════════════════════════════════════════════════════
# АВТОРИЗАЦІЯ ВЧИТЕЛЯ
# ═══════════════════════════════════════════════════════════════════════════════

def teacher_login(request):
    """Сторінка входу для вчителя."""
    if get_teacher_or_none(request):
        return redirect('teacher_dashboard')

    form = TeacherLoginForm()

    if request.method == 'POST':
        form = TeacherLoginForm(request.POST)
        if form.is_valid():
            username = form.cleaned_data['username']
            password = form.cleaned_data['password']
            user = authenticate(request, username=username, password=password)

            if user is not None:
                # Перевіряємо що є профіль вчителя
                try:
                    teacher = user.teacher_profile
                    login(request, user)
                    messages.success(request, f'Вітаємо, {teacher.full_name}! 👋')
                    return redirect('teacher_dashboard')
                except Teacher.DoesNotExist:
                    messages.error(request, 'Профіль вчителя не знайдено для цього облікового запису.')
            else:
                messages.error(request, 'Неправильний логін або пароль.')

    return render(request, 'feed/teacher_login.html', {'form': form})


def teacher_logout(request):
    """Вихід вчителя."""
    logout(request)
    messages.info(request, 'Ви вийшли з системи.')
    return redirect('index')


# ═══════════════════════════════════════════════════════════════════════════════
# ПАНЕЛЬ ВЧИТЕЛЯ (захищена)
# ═══════════════════════════════════════════════════════════════════════════════

@teacher_required
def teacher_dashboard(request):
    """
    Головна панель вчителя.
    Відображає всі завдання, згруповані за статусом.
    Підтримує AJAX-запити для плавної анімації вкладок.
    """
    teacher = request.user.teacher_profile

    # Автоматично архівуємо завдання, опубліковані понад 14 днів тому
    auto_archive_expired_assignments()

    # Автоматично оновлюємо відкладені публікації
    Assignment.objects.filter(
        teacher=teacher,
        status=Assignment.STATUS_SCHEDULED,
        scheduled_at__lte=timezone.now()
    ).update(
        status=Assignment.STATUS_PUBLISHED,
        published_at=timezone.now()
    )

    # Фільтр по вкладці
    tab = request.GET.get('tab', 'published')
    status_map = {
        'published': Assignment.STATUS_PUBLISHED,
        'draft': Assignment.STATUS_DRAFT,
        'scheduled': Assignment.STATUS_SCHEDULED,
        'archived': Assignment.STATUS_ARCHIVED,
    }
    status_filter = status_map.get(tab, Assignment.STATUS_PUBLISHED)

    assignments = Assignment.objects.filter(
        teacher=teacher,
        status=status_filter
    ).prefetch_related('classes', 'files').select_related('subject').order_by('-created_at')

    # Підраховуємо кількість у кожній вкладці
    counts = {
        'published': Assignment.objects.filter(teacher=teacher, status=Assignment.STATUS_PUBLISHED).count(),
        'draft': Assignment.objects.filter(teacher=teacher, status=Assignment.STATUS_DRAFT).count(),
        'scheduled': Assignment.objects.filter(teacher=teacher, status=Assignment.STATUS_SCHEDULED).count(),
        'archived': Assignment.objects.filter(teacher=teacher, status=Assignment.STATUS_ARCHIVED).count(),
    }

    context = {
        'teacher': teacher,
        'assignments': assignments,
        'tab': tab,
        'counts': counts,
    }

    if request.headers.get('X-Requested-With') == 'XMLHttpRequest' or request.GET.get('fragment') == '1':
        return render(request, 'feed/teacher_dashboard_fragment.html', context)

    return render(request, 'feed/teacher_dashboard.html', context)


@teacher_required
def assignment_create(request):
    """Створення нового завдання."""
    teacher = request.user.teacher_profile
    form = AssignmentForm(teacher=teacher)

    if request.method == 'POST':
        form = AssignmentForm(teacher=teacher, data=request.POST, files=request.FILES)
        if form.is_valid():
            assignment = form.save_with_status(teacher=teacher)

            # Зберігаємо всі прикріплені файли
            for f in request.FILES.getlist('files'):
                AssignmentFile.objects.create(
                    assignment=assignment,
                    file=f,
                    original_name=f.name
                )

            # Зберігаємо додаткові посилання
            for url, label in zip(request.POST.getlist('extra_link_url'), request.POST.getlist('extra_link_label')):
                if url.strip():
                    AssignmentLink.objects.create(
                        assignment=assignment,
                        url=url.strip(),
                        label=label.strip()
                    )

            # Зберігаємо додаткові YouTube відео
            for yurl, ytitle in zip(request.POST.getlist('extra_youtube_url'), request.POST.getlist('extra_youtube_title')):
                if yurl.strip():
                    AssignmentYouTubeLink.objects.create(
                        assignment=assignment,
                        url=yurl.strip(),
                        title=ytitle.strip()
                    )

            status_labels = {
                Assignment.STATUS_PUBLISHED: 'опубліковано',
                Assignment.STATUS_DRAFT: 'збережено як чернетку',
                Assignment.STATUS_SCHEDULED: 'заплановано до публікації',
            }
            label = status_labels.get(assignment.status, 'збережено')
            messages.success(request, f'Завдання "{assignment.title}" — {label}! ✅')
            return redirect('teacher_dashboard')

    context = {
        'teacher': teacher,
        'form': form,
        'is_edit': False,
    }
    return render(request, 'feed/assignment_form.html', context)


@teacher_required
def assignment_edit(request, pk):
    """Редагування існуючого завдання."""
    teacher = request.user.teacher_profile
    assignment = get_object_or_404(Assignment, pk=pk, teacher=teacher)

    form = AssignmentForm(teacher=teacher, instance=assignment)
    existing_files = assignment.files.all()
    existing_links = assignment.additional_links.all()
    existing_youtube_links = assignment.youtube_links.all()

    if request.method == 'POST':
        form = AssignmentForm(
            teacher=teacher,
            data=request.POST,
            files=request.FILES,
            instance=assignment
        )

        if form.is_valid():
            # Видаляємо файли, відмічені для видалення
            files_to_delete = request.POST.getlist('delete_files')
            if files_to_delete:
                AssignmentFile.objects.filter(
                    id__in=files_to_delete,
                    assignment=assignment
                ).delete()

            # Зберігаємо оновлене завдання
            assignment = form.save_with_status(teacher=teacher)

            # Додаємо нові файли
            for f in request.FILES.getlist('files'):
                AssignmentFile.objects.create(
                    assignment=assignment,
                    file=f,
                    original_name=f.name
                )

            # Оновлюємо додаткові посилання (перезаписуємо)
            assignment.additional_links.all().delete()
            for url, label in zip(request.POST.getlist('extra_link_url'), request.POST.getlist('extra_link_label')):
                if url.strip():
                    AssignmentLink.objects.create(
                        assignment=assignment,
                        url=url.strip(),
                        label=label.strip()
                    )

            # Оновлюємо додаткові YouTube відео (перезаписуємо)
            assignment.youtube_links.all().delete()
            for yurl, ytitle in zip(request.POST.getlist('extra_youtube_url'), request.POST.getlist('extra_youtube_title')):
                if yurl.strip():
                    AssignmentYouTubeLink.objects.create(
                        assignment=assignment,
                        url=yurl.strip(),
                        title=ytitle.strip()
                    )

            messages.success(request, f'Завдання "{assignment.title}" оновлено! ✅')
            return redirect('teacher_dashboard')

    context = {
        'teacher': teacher,
        'form': form,
        'assignment': assignment,
        'existing_files': existing_files,
        'existing_links': existing_links,
        'existing_youtube_links': existing_youtube_links,
        'is_edit': True,
    }
    return render(request, 'feed/assignment_form.html', context)


@teacher_required
def assignment_delete(request, pk):
    """Видалення завдання (POST-запит)."""
    teacher = request.user.teacher_profile
    assignment = get_object_or_404(Assignment, pk=pk, teacher=teacher)

    if request.method == 'POST':
        title = assignment.title
        assignment.delete()
        messages.success(request, f'Завдання "{title}" видалено.')

    return redirect('teacher_dashboard')


@teacher_required
def assignment_unarchive(request, pk):
    """
    Розархівація завдання з архіву (POST-запит).
    Повертає статус 'published', зберігаючи первинну дату публікації 'published_at',
    та встановлює час розархівації у полі 'unarchived_at'.
    """
    teacher = request.user.teacher_profile
    assignment = get_object_or_404(Assignment, pk=pk)

    if not (request.user.is_superuser or assignment.teacher == teacher):
        messages.error(request, 'У вас немає прав для розархівації цього завдання.')
        return redirect('teacher_dashboard')

    if request.method == 'POST':
        assignment.status = Assignment.STATUS_PUBLISHED
        assignment.unarchived_at = timezone.now()
        # Первинна дата published_at залишається незмінною!
        assignment.save(update_fields=['status', 'unarchived_at', 'updated_at'])
        messages.success(request, f'Завдання «{assignment.title}» успішно розархівовано (первинну дату збережено)! ♻️')

    return redirect(f"{reverse('teacher_dashboard')}?tab=published")


@teacher_required
def assignment_duplicate(request, pk):
    """
    Дублювання завдання.
    Створює нову чернетку на основі обраного завдання.
    """
    teacher = request.user.teacher_profile
    original = get_object_or_404(Assignment, pk=pk, teacher=teacher)

    # Створюємо дублікат
    duplicate = Assignment.objects.create(
        teacher=teacher,
        subject=original.subject,
        title=f"[Копія] {original.title}",
        description=original.description,
        is_individual=original.is_individual,
        student_name=original.student_name,
        link_url=original.link_url,
        link_label=original.link_label,
        due_date=original.due_date,
        status=Assignment.STATUS_DRAFT,
        duplicated_from=original,
    )
    duplicate.classes.set(original.classes.all())

    # Копіюємо файли (посилання на ті самі файли, без фізичного копіювання)
    for f in original.files.all():
        AssignmentFile.objects.create(
            assignment=duplicate,
            file=f.file,
            original_name=f.original_name,
        )

    messages.success(request, f'Чернетку "{duplicate.title}" створено. Можна редагувати.')
    return redirect('assignment_edit', pk=duplicate.pk)


@teacher_required
def assignment_archive(request, pk):
    """Переміщення завдання до архіву."""
    teacher = request.user.teacher_profile
    assignment = get_object_or_404(Assignment, pk=pk, teacher=teacher)

    if request.method == 'POST':
        assignment.status = Assignment.STATUS_ARCHIVED
        assignment.save()
        messages.info(request, f'Завдання "{assignment.title}" переміщено до архіву.')

    return redirect('teacher_dashboard')


@teacher_required
def assignment_publish(request, pk):
    """Миттєва публікація чернетки або відкладеного завдання."""
    teacher = request.user.teacher_profile
    assignment = get_object_or_404(Assignment, pk=pk, teacher=teacher)

    if request.method == 'POST':
        assignment.status = Assignment.STATUS_PUBLISHED
        if not assignment.published_at:
            assignment.published_at = timezone.now()
        assignment.save()
        messages.success(request, f'Завдання "{assignment.title}" опубліковано! ✅')

    return redirect('teacher_dashboard')


@teacher_required
def teacher_profile(request):
    """
    Редагування профілю вчителя:
    - Персональні предмети та класи вчителя (додавання/видалення з власного профілю).
    - Створення нових предметів та класів.
    - Функції Супер-Адміністратора: створення нових вчителів та скидання їхніх паролів.
    """
    teacher = request.user.teacher_profile
    profile_form = TeacherProfileForm(instance=teacher)
    subject_form = SubjectForm()
    class_form = ClassGroupForm()
    teacher_create_form = TeacherCreateForm()
    password_reset_form = PasswordResetForm()

    if request.method == 'POST':
        action = request.POST.get('action')

        # ── Оновлення основних даних (ПІБ, колір, фото аватара) ──────────────
        if action == 'update_profile':
            profile_form = TeacherProfileForm(request.POST, request.FILES, instance=teacher)
            if profile_form.is_valid():
                profile_form.save()
                messages.success(request, 'Профіль та аватар успішно збережено! ✅')
                return redirect('teacher_profile')
            else:
                for field, errors in profile_form.errors.items():
                    for err in errors:
                        messages.error(request, f'{err}')

        elif action == 'remove_avatar':
            if teacher.avatar_image:
                try:
                    import os
                    if os.path.exists(teacher.avatar_image.path):
                        os.remove(teacher.avatar_image.path)
                except Exception:
                    pass
                teacher.avatar_image = None
                teacher.save(update_fields=['avatar_image'])
                messages.info(request, 'Фото аватара видалено. Використовуються кольорові ініціали.')
            return redirect('teacher_profile')

        # ── Персональні предмети вчителя ─────────────────────────────────────
        elif action == 'add_teacher_subject':
            subject_id = request.POST.get('subject_id')
            if subject_id:
                subject = get_object_or_404(Subject, pk=subject_id)
                teacher.subjects.add(subject)
                messages.success(request, f'Предмет «{subject.name}» додано до вашого профілю! 📚')
            return redirect('teacher_profile')

        elif action == 'remove_teacher_subject':
            subject_id = request.POST.get('subject_id')
            if subject_id:
                subject = get_object_or_404(Subject, pk=subject_id)
                teacher.subjects.remove(subject)
                messages.info(request, f'Предмет «{subject.name}» видалено з вашого профілю.')
            return redirect('teacher_profile')

        # ── Персональні класи вчителя ─────────────────────────────────────────
        elif action == 'add_teacher_class':
            class_id = request.POST.get('class_id')
            if class_id:
                cgroup = get_object_or_404(ClassGroup, pk=class_id)
                teacher.classes.add(cgroup)
                messages.success(request, f'Клас «{cgroup.name}» додано до вашого профілю! 🏫')
            return redirect('teacher_profile')

        elif action == 'remove_teacher_class':
            class_id = request.POST.get('class_id')
            if class_id:
                cgroup = get_object_or_404(ClassGroup, pk=class_id)
                teacher.classes.remove(cgroup)
                messages.info(request, f'Клас «{cgroup.name}» видалено з вашого профілю.')
            return redirect('teacher_profile')

        # ── Створення нового предмету (і авто-додавання вчителю) ──────────────
        elif action == 'create_subject':
            subject_form = SubjectForm(request.POST)
            if subject_form.is_valid():
                subject = subject_form.save(commit=False)
                subject.created_by = teacher
                subject.save()
                teacher.subjects.add(subject)
                messages.success(request, f'Новий предмет «{subject.name}» створено та додано до профілю! 📚')
                return redirect('teacher_profile')
            else:
                messages.error(request, 'Помилка при створенні предмету.')

        # ── Створення нового класу (і авто-додавання вчителю) ─────────────────
        elif action == 'create_class':
            class_form = ClassGroupForm(request.POST)
            if class_form.is_valid():
                cgroup = class_form.save(commit=False)
                cgroup.created_by = teacher
                cgroup.save()
                teacher.classes.add(cgroup)
                messages.success(request, f'Новий клас «{cgroup.name}» створено та додано до профілю! 🏫')
                return redirect('teacher_profile')
            else:
                messages.error(request, 'Помилка при створенні класу.')

        # ── Функції Супер-Адміністратора: Створення акаунта вчителя ─────────
        elif action == 'create_teacher_account':
            if not request.user.is_superuser:
                messages.error(request, 'Тільки Супер-Адміністратор може створювати акаунти.')
                return redirect('teacher_profile')

            teacher_create_form = TeacherCreateForm(request.POST)
            if teacher_create_form.is_valid():
                username = teacher_create_form.cleaned_data['username']
                password = teacher_create_form.cleaned_data['password']
                full_name = teacher_create_form.cleaned_data['full_name']
                is_super = teacher_create_form.cleaned_data['is_superuser']

                new_user = User.objects.create_user(username=username, password=password)
                new_user.is_superuser = is_super
                new_user.is_staff = is_super
                new_user.save()

                Teacher.objects.create(user=new_user, full_name=full_name)
                role = "Супер-Адміністратор" if is_super else "Вчитель"
                messages.success(request, f'Акаунт ({role}) «{full_name}» ({username}) створено! 👤')
                return redirect('teacher_profile')
            else:
                for field, errors in teacher_create_form.errors.items():
                    for err in errors:
                        messages.error(request, f'{err}')

        # ── Функції Супер-Адміністратора: Скидання пароля ────────────────────
        elif action == 'reset_teacher_password':
            if not request.user.is_superuser:
                messages.error(request, 'Тільки Супер-Адміністратор може скидати паролі.')
                return redirect('teacher_profile')

            target_user_id = request.POST.get('user_id')
            new_pass = request.POST.get('new_password', '').strip()

            if target_user_id and new_pass:
                target_user = get_object_or_404(User, pk=target_user_id)
                target_user.set_password(new_pass)
                target_user.save()
                messages.success(request, f'Пароль для користувача {target_user.username} змінено! 🔑')
            else:
                messages.error(request, 'Вкажіть новий пароль.')
            return redirect('teacher_profile')

        # ── Функції Супер-Адміністратора: Видалення вчителя та матеріалів ─────
        elif action == 'delete_teacher':
            if not request.user.is_superuser:
                messages.error(request, 'Тільки Супер-Адміністратор може видаляти вчителів.')
                return redirect('teacher_profile')

            target_teacher_id = request.POST.get('teacher_id')
            if target_teacher_id:
                target_teacher = get_object_or_404(Teacher, pk=target_teacher_id)
                
                # Запобігаємо видаленню самого себе
                if target_teacher.user == request.user:
                    messages.error(request, 'Ви не можете видалити власний акаунт.')
                    return redirect('teacher_profile')

                # Видаляємо фізичні файли завдань вчителя з диска
                import os
                for assignment in target_teacher.assignments.all():
                    for file_obj in assignment.files.all():
                        if file_obj.file and os.path.exists(file_obj.file.path):
                            try:
                                os.remove(file_obj.file.path)
                            except Exception:
                                pass

                # Видаляємо користувача (і пов'язаного вчителя cascade)
                user_to_delete = target_teacher.user
                full_name_deleted = target_teacher.full_name
                username_deleted = user_to_delete.username
                user_to_delete.delete()

                messages.success(request, f'Вчителя «{full_name_deleted}» ({username_deleted}) та всі його матеріали успішно видалено! 🗑️')
            else:
                messages.error(request, 'Вкажіть вчителя для видалення.')
            return redirect('teacher_profile')


    # Отримуємо дані для сторінки (тільки предмети та класи, створені вчителями)
    teacher_subjects = teacher.subjects.all()
    teacher_classes = teacher.classes.all().order_by('grade', 'letter')

    other_subjects = Subject.objects.filter(created_by__isnull=False).exclude(id__in=teacher_subjects.values_list('id', flat=True))
    other_classes = ClassGroup.objects.filter(created_by__isnull=False).exclude(id__in=teacher_classes.values_list('id', flat=True)).order_by('grade', 'letter')


    all_teachers = Teacher.objects.select_related('user').all() if request.user.is_superuser else []

    context = {
        'teacher': teacher,
        'form': profile_form,
        'subject_form': subject_form,
        'class_form': class_form,
        'teacher_create_form': teacher_create_form,
        'password_reset_form': password_reset_form,
        'teacher_subjects': teacher_subjects,
        'teacher_classes': teacher_classes,
        'other_subjects': other_subjects,
        'other_classes': other_classes,
        'all_teachers': all_teachers,
    }
    return render(request, 'feed/teacher_profile.html', context)


from django.views.decorators.clickjacking import xframe_options_exempt


def get_pdf_preview_url(file_obj):
    """
    Конвертує pptx, docx, або xlsx у PDF за допомогою headless LibreOffice, якщо його ще немає в кеші.
    Повертає URL до PDF файлу або None у разі помилки.
    """
    import os
    import subprocess
    from django.conf import settings
    
    previews_dir = os.path.join(settings.MEDIA_ROOT, 'previews')
    os.makedirs(previews_dir, exist_ok=True)
    
    pdf_filename = f"{file_obj.id}.pdf"
    pdf_path = os.path.join(previews_dir, pdf_filename)
    
    if os.path.exists(pdf_path):
        from django.urls import reverse
        return reverse('file_view', args=[file_obj.id]) + "?preview_pdf=1"
        
    try:
        # Запускаємо LibreOffice для конвертації
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', previews_dir,
            file_obj.file.path
        ]
        subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True, timeout=20)
        
        # Отримуємо назву оригінального файлу та міняємо розширення на .pdf
        base_name = os.path.basename(file_obj.file.path)
        raw_pdf_name = os.path.splitext(base_name)[0] + '.pdf'
        raw_pdf_path = os.path.join(previews_dir, raw_pdf_name)
        
        if os.path.exists(raw_pdf_path):
            os.rename(raw_pdf_path, pdf_path)
            from django.urls import reverse
            return reverse('file_view', args=[file_obj.id]) + "?preview_pdf=1"
    except Exception as e:
        print(f"Error converting file {file_obj.id} to PDF preview: {str(e)}")
        
    return None


@xframe_options_exempt
def file_view(request, file_id):
    """Служить файл безпосередньо у браузері з інлайновим Content-Disposition та правильним MIME-типом (зокрема для PDF)."""
    file_obj = get_object_or_404(AssignmentFile, pk=file_id)
    try:
        import mimetypes
        import os
        from django.conf import settings
        
        # Якщо запит йде на PDF прев'ю сконвертованого офісного файлу
        if request.GET.get('preview_pdf') == '1':
            path = os.path.join(settings.MEDIA_ROOT, 'previews', f"{file_obj.id}.pdf")
            mime_type = 'application/pdf'
        else:
            path = file_obj.file.path
            mime_type, _ = mimetypes.guess_type(path)
            if not mime_type:
                mime_type = 'application/octet-stream'
                
        if os.path.exists(path):
            with open(path, 'rb') as f:
                response = HttpResponse(f.read(), content_type=mime_type)
            
            # Насильно змушуємо браузер рендерити inline (у фреймі)
            response['Content-Disposition'] = f'inline; filename="{file_obj.original_name}"'
            return response
    except Exception as e:
        return HttpResponse(f"Помилка завантаження файлу: {str(e)}", status=500)
    return HttpResponse("Файл не знайдено", status=404)


def file_preview(request, file_id):
    """Служба для генерації інлайнового прев'ю файлу (docx, xlsx, pptx, python та інших код-файлів)."""
    file_obj = get_object_or_404(AssignmentFile, pk=file_id)
    ext = file_obj.get_extension()
    
    # Спочатку пробуємо якісну конвертацію у PDF через LibreOffice
    if ext in {'.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls'}:
        pdf_url = get_pdf_preview_url(file_obj)
        if pdf_url:
            return JsonResponse({
                'type': 'url',
                'url': pdf_url,
                'file_type': 'pdf'
            })
            
    # Резервна обробка (fallback), якщо LibreOffice не зміг згенерувати PDF
    if ext == '.docx':
        try:
            import mammoth
            with open(file_obj.file.path, 'rb') as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html = result.value
                styled_html = f"<div class='docx-preview-content' style='text-align:left; width:100%; color:var(--color-text-primary); line-height:1.6;'>{html}</div>"
                return JsonResponse({
                    'type': 'html',
                    'content': styled_html
                })
        except Exception as e:
            return JsonResponse({
                'type': 'error',
                'message': f'Помилка конвертації файлу Word: {str(e)}'
            })

    elif ext == '.xlsx':
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_obj.file.path, read_only=True, data_only=True)
            html = []
            for sheet_name in wb.sheetnames[:3]:
                sheet = wb[sheet_name]
                html.append(f"<h3 style='margin-top:20px; margin-bottom:10px; color:var(--color-primary); text-align:left;'>📊 Аркуш: {sheet_name}</h3>")
                html.append("<div style='overflow-x:auto; width:100%; border-radius:8px; border:1px solid var(--color-border); margin-bottom:20px;'><table style='width:100%; border-collapse:collapse; font-size:13px; background:var(--color-surface);'>")
                
                for r_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                    if r_idx > 100:
                        html.append("<tr><td colspan='100' style='text-align:center; color:var(--color-text-muted); padding:8px;'>... відображено перші 100 рядків ...</td></tr>")
                        break
                    if not any(row):
                        continue
                        
                    html.append("<tr style='border-bottom:1px solid var(--color-border);'>")
                    for cell_value in row:
                        val = str(cell_value) if cell_value is not None else ""
                        if r_idx == 0:
                            html.append(f"<th style='border-right:1px solid var(--color-border); padding:8px; background:var(--color-bg-secondary); font-weight:600; text-align:left;'>{val}</th>")
                        else:
                            html.append(f"<td style='border-right:1px solid var(--color-border); padding:8px; text-align:left;'>{val}</td>")
                    html.append("</tr>")
                html.append("</table></div>")
            
            styled_html = "".join(html)
            return JsonResponse({
                'type': 'html',
                'content': styled_html
            })
        except Exception as e:
            return JsonResponse({
                'type': 'error',
                'message': f'Помилка конвертації таблиці Excel: {str(e)}'
            })

    elif ext == '.pptx':
        try:
            import base64
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.util import Emu

            prs = Presentation(file_obj.file.path)
            slide_w = prs.slide_width or Emu(9144000)
            slide_h = prs.slide_height or Emu(5143500)
            aspect = (slide_h / slide_w * 100) if slide_w else 56.25

            html = ["<div class='pptx-preview-wrapper' style='display:flex; flex-direction:column; gap:24px; width:100%;'>"]

            for idx, slide in enumerate(prs.slides):
                # ─── Background colour / gradient ───────────────────────────────
                bg_css = "background: #f0f0f0;"
                try:
                    bg = slide.background
                    fill = bg.fill
                    fill_type = fill.type
                    if fill_type is not None:
                        if hasattr(fill, 'fore_color') and fill.fore_color.type is not None:
                            try:
                                rgb = fill.fore_color.rgb
                                bg_css = f"background: #{rgb};"
                            except Exception:
                                pass
                except Exception:
                    pass

                # ─── Slide container ─────────────────────────────────────────────
                html.append(
                    f"<div class='pptx-slide' style='position:relative; width:100%; padding-bottom:{aspect:.2f}%; "
                    f"border-radius:8px; {bg_css} box-shadow:0 2px 12px rgba(0,0,0,.15); overflow:hidden;'>"
                )
                html.append("<div style='position:absolute; top:0; left:0; width:100%; height:100%; overflow:hidden;'>")
                html.append(f"<div style='position:absolute; top:6px; left:8px; font-size:10px; opacity:.5; z-index:10;'>Слайд {idx+1}</div>")

                # ─── Shapes ───────────────────────────────────────────────────────
                shapes = sorted(slide.shapes, key=lambda s: (s.top or 0))
                for shape in shapes:
                    try:
                        left_pct  = (shape.left  or 0) / slide_w * 100
                        top_pct   = (shape.top   or 0) / slide_h * 100
                        width_pct = (shape.width or 0) / slide_w * 100
                        height_pct= (shape.height or 0) / slide_h * 100
                    except Exception:
                        continue

                    pos_style = (
                        f"position:absolute; left:{left_pct:.2f}%; top:{top_pct:.2f}%; "
                        f"width:{width_pct:.2f}%; min-height:{height_pct:.2f}%; overflow:hidden; box-sizing:border-box;"
                    )

                    # -- Images / pictures --
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img_bytes = shape.image.blob
                            img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                            mime = shape.image.content_type or 'image/png'
                            html.append(
                                f"<img src='data:{mime};base64,{img_b64}' "
                                f"style='{pos_style} object-fit:contain; max-width:100%; max-height:100%;' alt=''>"
                            )
                        except Exception:
                            pass
                        continue

                    # -- Text frames --
                    if shape.has_text_frame:
                        tf = shape.text_frame
                        text_css = pos_style + " display:flex; flex-direction:column; justify-content:center; padding:2%;"

                        # Try to detect fill colour for text boxes
                        try:
                            sfill = shape.fill
                            if sfill.type is not None:
                                srgb = sfill.fore_color.rgb
                                text_css += f" background:#{srgb}CC;"
                        except Exception:
                            pass

                        is_title = (shape.is_placeholder and
                                    hasattr(shape, 'placeholder_format') and
                                    shape.placeholder_format.idx == 0)

                        inner_html = []
                        for para in tf.paragraphs:
                            para_text = para.text.strip()
                            if not para_text:
                                continue
                            colour = "#333"
                            font_size = "3"
                            bold = False
                            try:
                                run = para.runs[0] if para.runs else None
                                if run:
                                    if run.font.size:
                                        pts = run.font.size.pt
                                        font_size = f"{max(1, min(int(pts * 0.6), 8))}"
                                    if run.font.color.type is not None:
                                        colour = f"#{run.font.color.rgb}"
                                    bold = run.font.bold or False
                            except Exception:
                                pass

                            bw = "700" if (is_title or bold) else "400"
                            inner_html.append(
                                f"<p style='margin:0 0 2px; font-size:{font_size}cqw; font-weight:{bw}; "
                                f"color:{colour}; line-height:1.3; white-space:pre-wrap;'>{para_text}</p>"
                            )

                        if inner_html:
                            html.append(f"<div style='{text_css} container-type:inline-size;'>{''.join(inner_html)}</div>")

                html.append("</div>")  # inner absolute
                html.append("</div>")  # slide container

            html.append("</div>")  # wrapper
            return JsonResponse({'type': 'html', 'content': ''.join(html)})
        except Exception as e:
            return JsonResponse({
                'type': 'error',
                'message': f'Помилка конвертації презентації PowerPoint: {str(e)}'
            })

    # 4. Обробка текстових файлів та код-файлів (включаючи .py, .js, .css, .json, .sh тощо)
    elif file_obj.get_file_type() == 'text' or ext in {'.py', '.js', '.css', '.json', '.sh', '.cpp', '.h', '.c', '.java'}:
        try:
            import os
            if os.path.exists(file_obj.file.path):
                with open(file_obj.file.path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return JsonResponse({
                    'type': 'text',
                    'content': content,
                    'lang': ext[1:]
                })
        except Exception as e:
            return JsonResponse({
                'type': 'error',
                'message': f'Не вдалося прочитати текстовий файл: {str(e)}'
            })
            
    # 5. Дефолтна обробка для інших типів (зображення, pdf, відео, аудіо) -> перенаправляємо на inline-view
    from django.urls import reverse
    inline_view_url = reverse('file_view', args=[file_obj.id])
    return JsonResponse({
        'type': 'url',
        'url': inline_view_url,
        'file_type': file_obj.get_file_type()
    })


def server_stats(request):
    """
    API ендпоінт для PyQt6 лаунчера та моніторингу:
    Повертає статистику сервера та кількість підключених онлайн-клієнтів.
    """
    from .middleware import get_online_stats
    stats = get_online_stats()
    stats['status'] = 'running'
    stats['total_assignments'] = Assignment.objects.filter(status=Assignment.STATUS_PUBLISHED).count()
    return JsonResponse(stats)


from django.views.decorators.csrf import csrf_exempt


def client_heartbeat(request):
    """
    Періодичний heartbeat від браузера на всіх сторінках сайту.
    Оновлює активність та поточну сторінку клієнта.
    """
    import time
    from datetime import datetime
    from .middleware import get_client_ip, _active_clients, _lock
    
    ip = get_client_ip(request)
    current_path = request.GET.get('path', '/')
    now = time.time()
    with _lock:
        _active_clients[ip] = {
            'timestamp': now,
            'last_path': current_path,
            'last_time_str': datetime.now().strftime('%H:%M:%S')
        }
    return HttpResponse("OK")


@csrf_exempt
def client_disconnect(request):
    """
    Викликається браузером через navigator.sendBeacon при закритті вкладки/сайту,
    щоб миттєво відняти клієнта з лічильника онлайн.
    """
    from .middleware import get_client_ip, remove_client_ip
    ip = get_client_ip(request)
    remove_client_ip(ip)
    return HttpResponse("OK")



